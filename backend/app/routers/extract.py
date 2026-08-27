"""Extract router: terima file, jalankan extractor, return code blocks.

Mendukung:
- PDF: via pdf_extract.py sidecar (font-analysis + OCR fallback)
- MD, IPYNB, TXT, TEX, HTML: parse langsung di Python (tidak butuh sidecar)
- DOCX, PPTX, XLSX: TODO di iterasi berikutnya (bisa pakai python-docx, python-pptx)

Deteksi bahasa dilakukan di BE via language_detection.py (pygments + custom override).
"""
import os
import sys
import json
import re
import tempfile
import subprocess
from pathlib import Path
from fastapi import APIRouter, UploadFile, File, HTTPException, status, Request
from pydantic import BaseModel

from ..language_detection import detect_language
from ..pattern_extract import extract_from_pdf as pattern_extract_pdf, extract_from_text as pattern_extract_text
from ..rate_limit import limiter
from ..cache import get_cached, set_cached, get_cache_stats, clear_cache

router = APIRouter(prefix="/extract", tags=["extract"])


# Path ke script Python (relative ke file ini: backend/app/routers/extract.py)
# parents[0] = routers, parents[1] = app, parents[2] = backend
SIDECAR_SCRIPT = Path(__file__).resolve().parents[2] / "scripts" / "pdf_extract.py"
SIDECAR_TIMEOUT = 300  # 5 menit (OCR butuh waktu lama)

# Format yang bisa diproses langsung di BE (tanpa sidecar)
PLAIN_TEXT_EXTS = {"txt", "tex", "latex", "sty", "cls"}
MARKDOWN_EXTS = {"md", "markdown"}
IPYNB_EXTS = {"ipynb"}
HTML_EXTS = {"html", "htm"}

# Format yang butuh sidecar Python pdf_extract.py
PDF_EXTS = {"pdf"}

# Format yang didukung oleh officeparser (Node.js) — untuk V1, masih pakai FE route lama
# TODO: integrasikan python-docx/python-pptx di iterasi berikutnya
OFFICE_EXTS = {"docx", "pptx", "xlsx"}

ALL_SUPPORTED_EXTS = (
    PDF_EXTS | PLAIN_TEXT_EXTS | MARKDOWN_EXTS | IPYNB_EXTS | HTML_EXTS | OFFICE_EXTS
)


class CodeBlock(BaseModel):
    index: int
    lang: str
    code: str
    lines: int
    source: str = "font"


class ExtractResponse(BaseModel):
    blocks: list[CodeBlock]
    filename: str
    size: int
    total: int
    stats: dict | None = None


@router.post("", response_model=ExtractResponse)
@limiter.limit("10/hour")
async def extract_code(request: Request, file: UploadFile = File(...)):
    """Ekstrak code blocks dari file yang di-upload.

    Format didukung:
    - PDF (font-based + OCR fallback via Tesseract)
    - Markdown (.md) — fenced code blocks
    - IPYNB (.ipynb) — code cells
    - HTML (.html) — text content dengan code blocks
    - TXT / TEX / LATEX — plain text + LaTeX verbatim/lstlisting
    - DOCX, PPTX, XLSX — via python-docx / python-pptx / openpyxl

    File TIDAK disimpan ke DB. Hanya filename yang dipreserve di response.
    Hasil di-cache 24 jam berdasarkan hash file content (Redis kalau tersedia,
    in-memory fallback).
    """
    # Baca file
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="File kosong")

    # Cek ukuran
    max_bytes = int(os.environ.get("CL_MAX_UPLOAD_MB", "50")) * 1024 * 1024
    if len(content) > max_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"Ukuran file {len(content)/1024/1024:.1f}MB melebihi batas {max_bytes/1024/1024:.0f}MB"
        )

    # Validasi extension
    filename = file.filename or "upload.pdf"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ALL_SUPPORTED_EXTS:
        raise HTTPException(
            status_code=400,
            detail=f"Format .{ext} tidak didukung. Format yang didukung: {', '.join(sorted(ALL_SUPPORTED_EXTS))}"
        )

    # ─── Cek cache dulu ───
    # Kalau file yang sama sudah pernah di-extract, return dari cache
    # supaya tidak perlu jalankan sidecar lagi (hemat 5-60 detik).
    # TAPI: jangan return cache kalau result-nya empty (0 blocks) — kemungkinan
    # hasil dari versi BE lama yang belum punya heuristic fallback.
    # Bypass cache kalau empty, supaya selalu coba extract ulang.
    cached = get_cached(content)
    if cached is not None and len(cached.get("blocks", [])) > 0:
        # Update filename (mungkin user upload nama beda tapi content sama)
        cached["filename"] = filename
        cached["cached"] = True
        return ExtractResponse(**cached)

    # Route ke extractor yang sesuai
    if ext in PDF_EXTS:
        # PDF: pattern-based extraction (marker + density)
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            raw_blocks = pattern_extract_pdf(tmp_path)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass
        blocks = [
            CodeBlock(
                index=i,
                lang=b.get("lang", "unknown"),
                code=b["code"],
                lines=b["lines"],
                source=b.get("source", "pattern"),
            )
            for i, b in enumerate(raw_blocks)
        ]
    elif ext in MARKDOWN_EXTS:
        raw_blocks = pattern_extract_text(content.decode("utf-8", errors="replace"))
        blocks = [
            CodeBlock(index=i, lang=b.get("lang", "unknown"), code=b["code"],
                      lines=b["lines"], source=b.get("source", "pattern"))
            for i, b in enumerate(raw_blocks)
        ]
    elif ext in IPYNB_EXTS:
        blocks = extract_ipynb(content.decode("utf-8", errors="replace"))
    elif ext in HTML_EXTS:
        blocks = extract_html(content.decode("utf-8", errors="replace"))
    elif ext in PLAIN_TEXT_EXTS:
        raw_blocks = pattern_extract_text(content.decode("utf-8", errors="replace"))
        blocks = [
            CodeBlock(index=i, lang=b.get("lang", "unknown"), code=b["code"],
                      lines=b["lines"], source=b.get("source", "pattern"))
            for i, b in enumerate(raw_blocks)
        ]
    elif ext in OFFICE_EXTS:
        blocks = extract_office(content, ext)
    else:
        raise HTTPException(status_code=400, detail=f"Format .{ext} tidak didukung")

    response = ExtractResponse(
        blocks=blocks,
        filename=filename,
        size=len(content),
        total=len(blocks),
        stats=None,
    )

    # Simpan ke cache untuk request berikutnya
    set_cached(content, response.model_dump())

    return response


@router.get("/cache/stats")
def cache_stats():
    """Stats cache untuk debugging/admin."""
    return get_cache_stats()


@router.delete("/cache")
def cache_clear():
    """Hapus semua cache (admin only — TODO: tambah auth admin)."""
    count = clear_cache()
    return {"cleared": count}


# ─── PDF extraction via sidecar ───
async def extract_pdf(content: bytes, ext: str) -> list[CodeBlock]:
    """Ekstrak code blocks dari PDF via pdf_extract.py sidecar."""
    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        proc = subprocess.run(
            [sys.executable, str(SIDECAR_SCRIPT), tmp_path],
            capture_output=True, text=True, timeout=SIDECAR_TIMEOUT
        )
    except subprocess.TimeoutExpired:
        raise HTTPException(
            status_code=504,
            detail=f"Proses ekstraksi timeout setelah {SIDECAR_TIMEOUT}s. File mungkin terlalu besar atau berbasis gambar (OCR lambat)."
        )
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    if proc.returncode != 0:
        raise HTTPException(
            status_code=500,
            detail=f"Sidecar error: {proc.stderr[:500]}"
        )

    try:
        # Robust JSON parse: kalau ada warning text sebelum JSON (mis. dari
        # PyMuPDF deprecation warning yang bocor ke stdout), find first '{'
        # dan parse dari situ.
        stdout = proc.stdout.strip()
        json_start = stdout.find("{")
        if json_start > 0:
            # Ada prefix sebelum JSON — skip (kemungkinan warning text)
            stdout = stdout[json_start:]
        elif json_start == -1:
            # Tidak ada JSON sama sekali
            raise json.JSONDecodeError("No JSON found in output", stdout, 0)
        result = json.loads(stdout)
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=500,
            detail=f"Sidecar output bukan JSON valid: {e}. Output: {proc.stdout[:300]}"
        )

    if "error" in result and "blocks" not in result:
        raise HTTPException(status_code=422, detail=result["error"])

    return [
        CodeBlock(
            index=i,
            lang=b.get("lang", "unknown"),  # akan di-override oleh detect_language
            code=b["code"],
            lines=b["lines"],
            source=b.get("source", "font"),
        )
        for i, b in enumerate(result.get("blocks", []))
    ]


# ─── Markdown extraction ───
FENCED_RE = re.compile(
    r"```(\w*)\n?([\s\S]*?)```|~~~(\w*)\n?([\s\S]*?)~~~",
    re.MULTILINE,
)


def extract_markdown(text: str) -> list[CodeBlock]:
    """Ekstrak fenced code blocks dari Markdown."""
    blocks = []
    idx = 0
    for m in FENCED_RE.finditer(text):
        hint = (m.group(1) or m.group(3) or "").lower().strip()
        code = (m.group(2) or m.group(4) or "").strip()
        if len(code) < 10:
            continue
        # Pakai hint kalau ada, kalau tidak akan di-detect nanti
        lang = hint if hint else "unknown"
        blocks.append(CodeBlock(
            index=idx,
            lang=lang,
            code=code,
            lines=code.count("\n") + 1,
            source="fenced",
        ))
        idx += 1
    return blocks


# ─── IPYNB extraction ───
def extract_ipynb(text: str) -> list[CodeBlock]:
    """Ekstrak code cells dari Jupyter Notebook."""
    try:
        nb = json.loads(text)
    except json.JSONDecodeError:
        raise HTTPException(status_code=422, detail="IPYNB file bukan JSON valid")

    blocks = []
    idx = 0
    for cell in nb.get("cells", []):
        if cell.get("cell_type") != "code":
            continue
        source = cell.get("source", "")
        if isinstance(source, list):
            code = "".join(source)
        else:
            code = str(source)
        code = code.strip()
        if len(code) < 10:
            continue
        # IPYNB default Python, biar detect_language confirm
        blocks.append(CodeBlock(
            index=idx,
            lang="unknown",  # akan di-detect nanti
            code=code,
            lines=code.count("\n") + 1,
            source="ipynb",
        ))
        idx += 1
    return blocks


# ─── HTML extraction ───
def extract_html(text: str) -> list[CodeBlock]:
    """Ekstrak <pre><code> blocks dari HTML."""
    # Cari <pre><code>...</code></pre> atau <code class="language-X">
    pre_code_re = re.compile(
        r"<pre[^>]*>\s*<code[^>]*>([\s\S]*?)</code>\s*</pre>",
        re.IGNORECASE,
    )
    code_class_re = re.compile(
        r'<code[^>]*class=["\']language-(\w+)["\'][^>]*>([\s\S]*?)</code>',
        re.IGNORECASE,
    )

    blocks = []
    idx = 0

    # Pattern 1: <pre><code class="language-X">...</code></pre>
    for m in pre_code_re.finditer(text):
        # Cek apakah ada class language-X
        full_match = m.group(0)
        class_match = code_class_re.search(full_match)
        if class_match:
            hint = class_match.group(1).lower()
            code = class_match.group(2)
        else:
            hint = ""
            code = m.group(1)

        # Unescape HTML entities
        code = (code
                .replace("&lt;", "<")
                .replace("&gt;", ">")
                .replace("&amp;", "&")
                .replace("&quot;", '"')
                .replace("&#39;", "'")
                .strip())
        if len(code) < 10:
            continue
        blocks.append(CodeBlock(
            index=idx,
            lang=hint if hint else "unknown",
            code=code,
            lines=code.count("\n") + 1,
            source="html",
        ))
        idx += 1

    return blocks


# ─── Plain text / LaTeX extraction ───
def extract_plain_text(text: str, ext: str) -> list[CodeBlock]:
    """Ekstrak code blocks dari plain text atau LaTeX file.

    Untuk .tex: cari \begin{lstlisting}, \begin{verbatim}, \begin{minted}
    Untuk .txt: pakai heuristic token-density (sederhana)
    """
    if ext in {"tex", "latex", "sty", "cls"}:
        return extract_latex(text)
    else:
        # .txt: heuristic sederhana
        return extract_txt_heuristic(text)


def extract_latex(text: str) -> list[CodeBlock]:
    """Ekstrak code dari LaTeX lstlisting/verbatim/minted."""
    blocks = []
    idx = 0

    # \begin{lstlisting}[language=X]...\end{lstlisting}
    lst_re = re.compile(
        r"\\begin\{lstlisting\}(?:\[language=([^\]]+)\])?([\s\S]*?)\\end\{lstlisting\}",
        re.MULTILINE,
    )
    for m in lst_re.finditer(text):
        hint = (m.group(1) or "").lower().strip()
        code = (m.group(2) or "").strip()
        if len(code) < 10:
            continue
        blocks.append(CodeBlock(
            index=idx,
            lang=hint if hint else "unknown",
            code=code,
            lines=code.count("\n") + 1,
            source="latex",
        ))
        idx += 1

    # \begin{verbatim}...\end{verbatim}
    verb_re = re.compile(
        r"\\begin\{verbatim\}([\s\S]*?)\\end\{verbatim\}",
        re.MULTILINE,
    )
    for m in verb_re.finditer(text):
        code = (m.group(1) or "").strip()
        if len(code) < 10:
            continue
        blocks.append(CodeBlock(
            index=idx,
            lang="unknown",  # akan di-detect
            code=code,
            lines=code.count("\n") + 1,
            source="latex",
        ))
        idx += 1

    # \begin{minted}{lang}...\end{minted}
    mint_re = re.compile(
        r"\\begin\{minted\}(?:\{?\s*(\w+)\s*\}?)?([\s\S]*?)\\end\{minted\}",
        re.MULTILINE,
    )
    for m in mint_re.finditer(text):
        hint = (m.group(1) or "").lower().strip()
        code = (m.group(2) or "").strip()
        if len(code) < 10:
            continue
        blocks.append(CodeBlock(
            index=idx,
            lang=hint if hint else "unknown",
            code=code,
            lines=code.count("\n") + 1,
            source="latex",
        ))
        idx += 1

    return blocks


def extract_txt_heuristic(text: str) -> list[CodeBlock]:
    """Heuristic untuk TXT/DOCX — pakai scoring kode-like per baris.

    Sama dengan heuristic di pdf_extract.py:
    - R-specific keywords (+3)
    - Assignment <- (+3)
    - Pipe %>% (+2)
    - Function call, comment, string (+1)
    - Prose penalty (-2)
    - ## R output dianggap "soft" (boleh ikut blok, tidak membentuk blok sendiri)
    - Block start detection: library(), data.frame(), # Kasus N
    - Post-process: strip ## R output, pisah blok yang mengandung ##
    """
    lines = text.split("\n")
    if not lines:
        return []

    def score_line(line: str) -> int:
        t = line.strip()
        if not t:
            return 0
        score = 0
        # R-specific keywords (bobot tinggi)
        if re.search(r"\b(library|require|data\.frame|read\.csv|read\.table|read\.xlsx|"
                     r"summary|lm|glm|aov|cor\.test|chisq\.test|t\.test|"
                     r"ggplot|plot|abline|hist|boxplot|"
                     r"qt|qnorm|qf|qchisq|pt|pnorm|"
                     r"sample|set\.seed|c\s*\(|seq|rep|"
                     r"head|tail|str|names|colnames|rownames|"
                     r"mean|median|sd|var|sum|sqrt|abs|round|"
                     r"cbind|rbind|merge|subset|transform|"
                     r"cat|paste|paste0|sprintf|"
                     r"def|class|import|from|return|if|else|elif|for|while|"
                     r"function|var|let|const|public|private|static|void|int|float|"
                     r"print|echo|SELECT|FROM|WHERE)\b", t):
            score += 3
        # R assignment operator
        if re.search(r"\w+\s*<-", t):
            score += 3
        if re.search(r"<-|->|%>%|:=", t):
            score += 2
        # Function call
        if re.search(r"\b\w+\s*\(", t):
            score += 1
        # Comment
        if re.match(r"^\s*#", t):
            score += 1
        # Code symbols
        if re.search(r"[(){}\[\];=<>+\-*/\\&|!?:,'\".]", t):
            score += 1
        # String literal
        if re.search(r'["\'].*["\']', t):
            score += 1
        # Prose penalty
        prose_words = re.findall(r"\b(?:dan|atau|yang|untuk|pada|dengan|dari|ke|di|ini|itu|"
                                r"adalah|akan|sebuah|seorang|mahasiswa|rata|selisih|"
                                r"proporsi|signifikan|berbeda|menggunakan|menghitung|"
                                r"the|and|or|for|with|from|to|in|of|a|an|is|are|was|were)\b",
                                t, re.IGNORECASE)
        if len(prose_words) >= 2:
            score -= 2
        return score

    def is_r_output(line):
        t = line.strip()
        return t.startswith("## ") or t.startswith("##\t") or t.startswith("[1] ")

    def is_block_start(line):
        t = line.strip()
        if re.match(r"^\s*library\s*\(", t):
            return True
        if re.search(r"<-\s*data\.frame\s*\(", t):
            return True
        if re.match(r"^\s*#\s*(kasus|soal|contoh|latihan)\s+\d", t, re.IGNORECASE):
            return True
        return False

    blocks = []
    current_block_lines = []
    idx = 0
    THRESHOLD = 2

    for line in lines:
        if score_line(line) >= THRESHOLD or is_r_output(line):
            if is_block_start(line) and current_block_lines:
                real_code_lines = [l for l in current_block_lines if not is_r_output(l)]
                if real_code_lines:
                    code = "\n".join(current_block_lines).strip()
                    if len(code) >= 10:
                        blocks.append(CodeBlock(
                            index=idx,
                            lang="unknown",
                            code=code,
                            lines=code.count("\n") + 1,
                            source="heuristic",
                        ))
                        idx += 1
                current_block_lines = []
            current_block_lines.append(line)
        else:
            if len(current_block_lines) >= 2:
                real_code_lines = [l for l in current_block_lines if not is_r_output(l)]
                if real_code_lines:
                    code = "\n".join(current_block_lines).strip()
                    if len(code) >= 10:
                        blocks.append(CodeBlock(
                            index=idx,
                            lang="unknown",
                            code=code,
                            lines=code.count("\n") + 1,
                            source="heuristic",
                        ))
                        idx += 1
            current_block_lines = []

    if len(current_block_lines) >= 2:
        real_code_lines = [l for l in current_block_lines if not is_r_output(l)]
        if real_code_lines:
            code = "\n".join(current_block_lines).strip()
            if len(code) >= 10:
                blocks.append(CodeBlock(
                    index=idx,
                    lang="unknown",
                    code=code,
                    lines=code.count("\n") + 1,
                    source="heuristic",
                ))

    # Post-process: strip ## R output, pisah blok yang mengandung ##
    final_blocks = []
    for b in blocks:
        code_lines = b.code.split("\n")
        r_output_indices = [i for i, l in enumerate(code_lines) if is_r_output(l)]
        if not r_output_indices:
            final_blocks.append(b)
            continue

        current_chunk = []
        for line in code_lines:
            if is_r_output(line):
                if current_chunk:
                    code = "\n".join(current_chunk).strip()
                    if len(code) >= 10:
                        final_blocks.append(CodeBlock(
                            index=len(final_blocks),
                            lang="unknown",
                            code=code,
                            lines=code.count("\n") + 1,
                            source="heuristic",
                        ))
                    current_chunk = []
            else:
                current_chunk.append(line)
        if current_chunk:
            code = "\n".join(current_chunk).strip()
            if len(code) >= 10:
                final_blocks.append(CodeBlock(
                    index=len(final_blocks),
                    lang="unknown",
                    code=code,
                    lines=code.count("\n") + 1,
                    source="heuristic",
                ))

    return final_blocks


# ─── Office format extraction (DOCX, PPTX, XLSX) ───
def extract_office(content: bytes, ext: str) -> list[CodeBlock]:
    """Ekstrak code blocks dari DOCX, PPTX, atau XLSX.

    Strategi:
    - DOCX: pakai python-docx. Iterasi paragraf, identifikasi baris kode via heuristic.
      Plus extract text dari tabel (sering berisi kode di modul praktikum).
    - PPTX: pakai python-pptx. Iterasi slides, extract text dari shapes.
    - XLSX: pakai openpyxl. Iterasi cells, gabung jadi text, lalu heuristic.
    """
    import io

    if ext == "docx":
        return extract_docx(content)
    elif ext == "pptx":
        return extract_pptx(content)
    elif ext == "xlsx":
        return extract_xlsx(content)
    else:
        raise HTTPException(status_code=400, detail=f"Format Office .{ext} tidak didukung")


def extract_docx(content: bytes) -> list[CodeBlock]:
    """Ekstrak dari DOCX via python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-docx tidak terinstall. Run: uv add python-docx"
        )

    import io
    doc = Document(io.BytesIO(content))

    # Gabung semua paragraf + tabel jadi satu text
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)

    # Extract dari tabel juga (sering berisi kode di modul praktikum)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    # Cell bisa multi-line, pisah
                    for line in text.split("\n"):
                        if line.strip():
                            lines.append(line)

    # Pakai heuristic yang sama dengan TXT untuk identifikasi kode
    text = "\n".join(lines)
    return extract_txt_heuristic(text)


def extract_pptx(content: bytes) -> list[CodeBlock]:
    """Ekstrak dari PPTX via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx tidak terinstall. Run: uv add python-pptx"
        )

    import io
    prs = Presentation(io.BytesIO(content))

    lines = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            lines.append(text)

    text = "\n".join(lines)
    return extract_txt_heuristic(text)


def extract_xlsx(content: bytes) -> list[CodeBlock]:
    """Ekstrak dari XLSX via openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="openpyxl tidak terinstall. Run: uv add openpyxl"
        )

    import io
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text = str(cell).strip()
                    if text:
                        lines.append(text)

    wb.close()
    text = "\n".join(lines)
    return extract_txt_heuristic(text)
